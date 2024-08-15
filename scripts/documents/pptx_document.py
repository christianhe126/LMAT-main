from documents.document import Document
import re
import base64
import sys
import json
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup

class PptxDocument(Document):
    def __init__(self, format, data, template=None, metadata=None):
        super().__init__(format, data, template, metadata)
        self.document = Presentation(self.path)
        if self.document.slides:
            self._add_title(self.document.slides[0], self.title)
        if self.metadata:
            self._add_metadata()
        self.max_lines = 15
        self.chars_per_line = 90
        self.paragraphs = []
        self.page_size = {
            "x": 242937,
            "y": 1093491,
            "width": 8656646,
            "height": 3683504
        }
    
    def create_page(self, title = "[No Title]", following = None):
        if self.document.slide_layouts[6]:
            page = self.document.slides.add_slide(self.document.slide_layouts[6])
        else:
            page = self.document.slides.add_slide()
        
        self._add_title(page, title)
        return page

    def _add_title(self, page, title):
        if "---" in title:
            parts = title.split("---")
            if len(parts) > 0:
                title = parts[0]

        if page.shapes.title:
            page.shapes.title.text = title
        else:
            textbox = page.shapes.add_textbox(242936, 323850, 6525307, 627720)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = title
            p.font.size = Pt(24)

    def add_image(self, page, x, y, width, height, source):
        b64_string = source
        if b64_string.startswith("data:image/"):
            b64_string = b64_string.split(",")[1]
        
        image_data = base64.b64decode(b64_string)
        image_stream = BytesIO(image_data)
        
        x, y, width, height = self._calculate_img_size_and_position(x, y, width, height, image_data)
        page.shapes.add_picture(image_stream, x + 100000, y + 100000, width - 200000, height - 200000)
    
    def add_text(self, page, x, y, width, height, source, font = None):
        text_frame = page.shapes.add_textbox(x,y,width,height).text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        self.html_to_element(text_frame, source, font = font)  

    def add_table(self, page, x, y, width, height, source):
        n_row, n_col, b_style, b_width, b_color, cellspacing, data, sizes = self._get_table_data(source)
        font = None
        if cellspacing and (cellspacing == 0 or cellspacing == "0"):
            font = 12

        if b_width == 0 or b_width == "0":
            regular_cell_width = int(width / n_col)
            regular_cell_height = int(height / n_row)
            for i in range(n_row):
                # Calculate the height
                _height = regular_cell_height
                if sizes and sizes[i][0][1]:  # Check if height is present
                    cell_height = float(sizes[i][0][1])
                    _height = Pt(cell_height * 0.75)
                _x = x
                for j in range(n_col):
                    if len(data) <= i or len(data[i]) <= j:
                        break

                    # Calculate the width
                    _width = regular_cell_width, 
                    if sizes and sizes[i][j][0]: # Check if width is present
                        cell_width_percentage = float(sizes[i][j][0])
                        _width = int((cell_width_percentage / 100) * self.page_size["width"])

                    # Add the content
                    if re.compile(r'<[^>]+>').search(data[i][j]):
                        soup = BeautifulSoup(data[i][j], 'html.parser')
                        if soup.find('img'):
                            self.add_image(page, _x, y, _width, _height, soup.find('img').get('src'))
                        elif soup.find('ul'):
                            uls = soup.findAll('ul')
                            self.add_text(page, _x, y, _width, _height, str(soup), font = font) #ul.decode_contents())
                            #self.html_to_element(text_frame, ul.decode_contents(), font = 12)
                        else:
                            self.add_text(page, _x, y, _width, _height, data[i][j], font = font)
                            #self.html_to_element(text_frame, data[i][j], font = 12)
                    else:
                        self.add_text(page, _x, y, _width, _height, data[i][j], font = font)
                        #self.html_to_element(text_frame, data[i][j], font = 12)
                    
                    _x = _x + _width

                y = y + _height

        else:
            shape = page.shapes.add_table(n_row, n_col,x,y,width,height)
            table = shape.table

            tbl = shape._element.graphic.graphicData.tbl
            style_id_transparent = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
            style_id_black_border = '{5940675A-B579-460E-94D1-54222C63F5DA}'
            style_id_theme_1 = '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'
            if b_width == 0 or b_width == "0":
                tbl[0][-1].text = style_id_transparent
            elif b_color:
                tbl[0][-1].text = style_id_theme_1
            else:
                tbl[0][-1].text = style_id_black_border
            
            for i in range(n_row):
                for j in range(n_col):
                    if len(data) <= i or len(data[i]) <= j:
                        break
                    cell = table.cell(i, j)
                    text_frame = cell.text_frame

                    # Set cell width and height
                    if sizes and sizes[i][j][0]: # Check if width is present
                        cell_width_percentage = float(sizes[i][j][0])
                        table.columns[j].width = int((cell_width_percentage / 100) * self.page_size["width"])
                    if sizes and sizes[i][j][1]:  # Check if height is present
                        cell_height = float(sizes[i][j][1])
                        table.rows[i].height = Pt(cell_height * 0.75)  # Convert pixel height to points (assuming 1px ≈ 0.75pt)

                    # Check if contains image
                    if re.compile(r'<[^>]+>').search(data[i][j]):
                        soup = BeautifulSoup(data[i][j], 'html.parser')
                        if soup.find('img'):
                            cell_x, cell_y, cell_width, cell_height = self._calculate_cell_position(table, i, j)
                            self.add_image(page, cell_x, cell_y, cell_width, cell_height, soup.find('img').get('src'))
                        elif soup.find('ul'):
                            ul = soup.find('ul')
                            self.html_to_element(text_frame, ul.decode_contents(), font = font)
                        else:
                            self.html_to_element(text_frame, data[i][j], font = font)
                    else:
                        self.html_to_element(text_frame, data[i][j], font = font)

    def _reduce_paragraphs(self, paragraphs):
        i = 0
        if len(paragraphs) == 1 and (paragraphs[0].type == "img" or paragraphs[0].type == "table"):
            paragraphs[0].estimate = self.max_lines - 1
            return

        while i < len(paragraphs) - 1:
            current = paragraphs[i]
            next_paragraph = paragraphs[i + 1]
            
            if current.type == "p" and next_paragraph.type == "p" and current.estimate + next_paragraph.estimate < self.max_lines:
                # Concatenate the contents and update the estimate
                current.content += next_paragraph.content
                current.estimate += next_paragraph.estimate
                # Remove the next paragraph from the list
                paragraphs.pop(i + 1)
            elif next_paragraph.type == "img" and next_paragraph.estimate < self.max_lines - current.estimate - 1:
                next_paragraph.estimate = self.max_lines - current.estimate - 1
                i += 1
            elif next_paragraph.type == "table" and next_paragraph.estimate < self.max_lines - current.estimate - 1:
                next_paragraph.estimate = self.max_lines - current.estimate - 1
                i += 1
            else:
                i += 1

    #def _format_bullet_point(self, text):
    #    return " " + text

    def _calculate_cell_position(self, table, row, col):
        table_gf = table._graphic_frame
        x = table_gf._element.xfrm.off.x
        y = table_gf._element.xfrm.off.y

        cell_left = x
        cell_top = y
        cell_height = table.rows[0].height
        cell_width = table.columns[0].width

        for r in range(row):
            cell_height = table.rows[r].height
            cell_top += cell_height

        for c in range(col):
            cell_width = table.columns[c].width
            cell_left += cell_width
        
        cell_height = table.rows[row].height
        cell_width = table.columns[col].width

        return cell_left, cell_top, cell_width, cell_height
    
    def _replace_text(self, old_text, new_text):
        for slide in self.document.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if(shape.text.find(old_text))!=-1:
                        text_frame = shape.text_frame
                        cur_text = text_frame.paragraphs[0].runs[0].text
                        _new_text = cur_text.replace(str(old_text), str(new_text))
                        text_frame.paragraphs[0].runs[0].text = _new_text

    def _add_metadata(self):
        for key, value in self.metadata.items():
            self._replace_text("###" + key + "###", value)
    
    def html_to_element(self, object, html, font = None):
        if object.paragraphs[0] and object.paragraphs[0].text == '':
            p = object.paragraphs[0]._p
            parent_element = p.getparent()
            parent_element.remove(p)
        
        super().html_to_element(object, html, font = font)
